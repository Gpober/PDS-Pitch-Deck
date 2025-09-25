#!/usr/bin/env python3
"""
Pride Dealer Services - HTML to PowerPoint Converter
Converts the HTML presentation to a professional PowerPoint deck
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup

def parse_html_file(html_path='index.html'):
    """Parse the HTML file and extract slide content"""
    
    if not os.path.exists(html_path):
        print(f"Error: {html_path} not found in current directory")
        return None
    
    with open(html_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    soup = BeautifulSoup(content, 'html.parser')
    slides_data = []
    
    # Find all slides
    slides = soup.find_all('div', class_='slide')
    
    for i, slide in enumerate(slides, 1):
        slide_data = {'slide_number': i}
        
        # Extract title
        title_elem = slide.find(['h1', 'h2'])
        if title_elem:
            slide_data['title'] = title_elem.get_text().strip()
        else:
            slide_data['title'] = f"Slide {i}"
        
        # Extract content parts
        content_parts = []
        
        # Get main paragraphs (excluding titles)
        paragraphs = slide.find_all('p')
        for p in paragraphs:
            text = p.get_text().strip()
            if text and text != slide_data.get('title', '') and len(text) > 10:
                content_parts.append(text)
        
        # Get metric boxes
        metric_boxes = slide.find_all('div', class_='metric-box')
        for box in metric_boxes:
            h3 = box.find('h3')
            if h3:
                content_parts.append(f"\n{h3.get_text().strip()}:")
            
            # Get metrics
            metric_values = box.find_all('span', class_='metric-value')
            metric_labels = box.find_all('span', class_='metric-label')
            
            for value, label in zip(metric_values, metric_labels):
                content_parts.append(f"• {value.get_text().strip()} - {label.get_text().strip()}")
            
            # Get paragraphs in metric box
            box_paras = box.find_all('p')
            for p in box_paras:
                text = p.get_text().strip()
                if text and len(text) > 10:
                    content_parts.append(f"• {text}")
        
        # Get key facts lists
        key_facts = slide.find_all('ul', class_='key-facts')
        for facts in key_facts:
            items = facts.find_all('li')
            for item in items:
                label_elem = item.find('span', class_='fact-label')
                value_elem = item.find('span', class_='fact-value')
                if label_elem and value_elem:
                    label = label_elem.get_text().strip()
                    value = value_elem.get_text().strip()
                    content_parts.append(f"• {label}: {value}")
                else:
                    # Regular list item
                    text = item.get_text().strip()
                    if text:
                        content_parts.append(f"• {text}")
        
        # Get regular lists
        regular_lists = slide.find_all('ul', class_=lambda x: x != 'key-facts' if x else True)
        for ul in regular_lists:
            if 'key-facts' not in (ul.get('class') or []):
                items = ul.find_all('li')
                for li in items:
                    text = li.get_text().strip()
                    if text and not text.startswith('•'):
                        content_parts.append(f"• {text}")
        
        # Get tables
        tables = slide.find_all('table', class_='financial-table')
        for table in tables:
            content_parts.append("\nFinancial Data:")
            headers = table.find_all('th')
            if headers:
                header_text = " | ".join([th.get_text().strip() for th in headers])
                content_parts.append(header_text)
                content_parts.append("-" * len(header_text))
            
            rows = table.find_all('tr')[1:]  # Skip header row
            for row in rows:
                cells = row.find_all('td')
                if cells:
                    row_text = " | ".join([cell.get_text().strip() for cell in cells])
                    content_parts.append(row_text)
        
        # Get partnership sections
        partnerships = slide.find_all('div', class_='partnership-section')
        for partnership in partnerships:
            h3 = partnership.find('h3')
            if h3:
                content_parts.append(f"\n{h3.get_text().strip()}:")
            
            facts = partnership.find_all('ul', class_='key-facts')
            for fact_list in facts:
                items = fact_list.find_all('li')
                for item in items:
                    label = item.find('span', class_='fact-label')
                    value = item.find('span', class_='fact-value')
                    if label and value:
                        content_parts.append(f"• {label.get_text().strip()}: {value.get_text().strip()}")
        
        slide_data['content'] = '\n'.join(content_parts)
        slides_data.append(slide_data)
        
        print(f"Extracted Slide {i}: {slide_data['title']}")
    
    return slides_data

def create_powerpoint_presentation(slides_data, output_filename='Pride_Dealer_Services_Presentation.pptx'):
    """Create PowerPoint presentation from extracted slide data"""
    
    # Create presentation
    prs = Presentation()
    
    # Define brand colors from your CSS
    brand_gold = RGBColor(212, 175, 55)    # #d4af37
    brand_dark = RGBColor(184, 148, 31)    # #b8941f  
    dark_blue = RGBColor(26, 26, 46)       # #1a1a2e
    surface_blue = RGBColor(22, 33, 62)    # #16213e
    accent_blue = RGBColor(15, 52, 96)     # #0f3460
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(204, 204, 204)
    
    def set_slide_background(slide, color):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def add_branded_slide(title_text, content_text, is_title_slide=False):
        """Add a slide with Pride Dealer Services branding"""
        
        if is_title_slide:
            slide_layout = prs.slide_layouts[0]  # Title slide
        else:
            slide_layout = prs.slide_layouts[1]  # Title and content
            
        slide = prs.slides.add_slide(slide_layout)
        
        # Set dark background
        set_slide_background(slide, dark_blue)
        
        if is_title_slide:
            # Title slide formatting
            title = slide.shapes.title
            subtitle = slide.placeholders[1]  # Subtitle placeholder
            
            title.text = title_text
            title_frame = title.text_frame
            title_frame.paragraphs[0].font.color.rgb = brand_gold
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            subtitle.text = "Investment Presentation\nNational Detail & Condition Reports Company"
            subtitle_frame = subtitle.text_frame
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.color.rgb = light_gray
                paragraph.font.size = Pt(18)
                paragraph.alignment = PP_ALIGN.CENTER
        else:
            # Regular slide formatting
            title = slide.shapes.title
            title.text = title_text
            title_frame = title.text_frame
            title_frame.paragraphs[0].font.color.rgb = brand_gold
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Content
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                content_frame = content_placeholder.text_frame
                content_frame.text = content_text
                
                # Format all paragraphs
                for paragraph in content_frame.paragraphs:
                    paragraph.font.color.rgb = white
                    paragraph.font.size = Pt(14)
                    paragraph.space_after = Pt(8)
                    
                    # Make bullet points stand out
                    if paragraph.text.strip().startswith('•'):
                        paragraph.font.size = Pt(12)
                        paragraph.space_before = Pt(4)
        
        return slide
    
    # Process all slides
    for i, slide_data in enumerate(slides_data):
        title = slide_data['title']
        content = slide_data['content']
        
        # Skip slides with no meaningful content
        if not content.strip() or len(content.strip()) < 20:
            print(f"Skipping slide {i+1} - insufficient content")
            continue
        
        # First slide is title slide
        is_title = (i == 0 and 'executive summary' in title.lower())
        
        # Clean up content
        content = content.replace('\n\n\n', '\n\n')  # Remove excessive line breaks
        content = content.strip()
        
        # Limit content length for readability
        if len(content) > 1500:
            content = content[:1500] + "..."
        
        add_branded_slide(title, content, is_title_slide=is_title)
        print(f"Created PowerPoint slide: {title}")
    
    # Save presentation
    prs.save(output_filename)
    print(f"\nPowerPoint presentation saved as: {output_filename}")
    print(f"Total slides created: {len(prs.slides)}")
    
    return output_filename

def main():
    """Main function to convert HTML to PowerPoint"""
    print("Pride Dealer Services - HTML to PowerPoint Converter")
    print("=" * 60)
    
    # Check if HTML file exists
    html_file = 'index.html'
    if not os.path.exists(html_file):
        print(f"Error: {html_file} not found in current directory")
        print("Make sure you're running this script in the same directory as your HTML file")
        return
    
    # Parse HTML file
    print(f"Parsing {html_file}...")
    slides_data = parse_html_file(html_file)
    
    if not slides_data:
        print("Failed to parse HTML file")
        return
    
    print(f"Found {len(slides_data)} slides")
    
    # Create PowerPoint presentation
    print("\nCreating PowerPoint presentation...")
    output_file = create_powerpoint_presentation(slides_data)
    
    print(f"\nConversion completed successfully!")
    print(f"Output file: {output_file}")

if __name__ == "__main__":
    main()
